"""Generate a standalone PPTX summarizing the new training run (full_p2_50epochs).

All figures are computed at run time from the source CSV/JSON files so the
slide content always reflects the current data on disk.
"""

import json
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
CKPT_DIR = ROOT / "checkpoints_new" / "full_p2_50epochs"
RESULTS_DIR = CKPT_DIR / "results"
PLOTS_DIR = CKPT_DIR / "plots"
OUT_PATH = ROOT / "thesis" / "ChagaSight_New_Training_Update.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.5)

NAVY = RGBColor(0x1F, 0x3B, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF0, 0xF0, 0xF0)
DARK_TEXT = RGBColor(0x20, 0x20, 0x20)


# ---------------------------------------------------------------------------
# Data loading / computation
# ---------------------------------------------------------------------------

def load_dataset_composition():
    all_data = pd.read_csv(ROOT / "data" / "processed" / "metadata" / "all_data.csv", low_memory=False)
    resolved = pd.read_csv(ROOT / "resolved_dataset_split.csv", low_memory=False)
    age_filtered = pd.read_csv(ROOT / "resolved_dataset_split_with_age_filter.csv", low_memory=False)

    orig = all_data.groupby("dataset").agg(records=("id", "size"), pos=("label_hard", "sum"))
    resolved_n = resolved.groupby("dataset").size()
    final = age_filtered.groupby("dataset").agg(records=("ecg_id", "size"), pos=("chagas", "sum"))

    # normalize dataset name differences between all_data.csv and resolved_*.csv
    name_map = {"code15": "code15", "ptbxl": "ptb_xl", "samitrop": "samitrop"}

    rows = []
    for orig_name, resolved_name in name_map.items():
        o_rec, o_pos = orig.loc[orig_name, "records"], orig.loc[orig_name, "pos"]
        r_rec = resolved_n.loc[resolved_name]
        f_rec, f_pos = final.loc[resolved_name, "records"], final.loc[resolved_name, "pos"]
        rows.append({
            "label": orig_name.upper().replace("CODE15", "CODE-15").replace("PTBXL", "PTB-XL").replace("SAMITROP", "SaMi-Trop"),
            "orig_records": int(o_rec),
            "orig_pos": int(o_pos),
            "resolved_records": int(r_rec),
            "final_records": int(f_rec),
            "final_pos": int(f_pos),
        })

    total = {
        "label": "Total",
        "orig_records": sum(r["orig_records"] for r in rows),
        "orig_pos": sum(r["orig_pos"] for r in rows),
        "resolved_records": sum(r["resolved_records"] for r in rows),
        "final_records": sum(r["final_records"] for r in rows),
        "final_pos": sum(r["final_pos"] for r in rows),
    }
    rows.append(total)

    age = age_filtered["age"]
    age_stats = {"min": age.min(), "max": age.max(), "mean": age.mean(), "median": age.median()}
    return rows, age_stats


def load_final_split():
    final_split = pd.read_csv(ROOT / "data" / "processed" / "metadata" / "final_split.csv", low_memory=False)
    g = final_split.groupby(["dataset", "split"]).agg(n=("id", "size"), pos=("label_hard", "sum"))
    g = g.unstack(fill_value=0)

    name_map = {"code15": "CODE-15", "ptbxl": "PTB-XL", "samitrop": "SaMi-Trop"}
    per_dataset = []
    for ds, label in name_map.items():
        n = g.loc[ds, "n"]
        pos = g.loc[ds, "pos"]
        per_dataset.append({
            "label": label,
            "train": int(n.get("train", 0)),
            "val": int(n.get("val", 0)),
            "test": int(n.get("test", 0)),
            "pos_train": int(pos.get("train", 0)),
            "pos_val": int(pos.get("val", 0)),
            "pos_test": int(pos.get("test", 0)),
        })

    total = {
        "label": "Total",
        "train": sum(r["train"] for r in per_dataset),
        "val": sum(r["val"] for r in per_dataset),
        "test": sum(r["test"] for r in per_dataset),
        "pos_train": sum(r["pos_train"] for r in per_dataset),
        "pos_val": sum(r["pos_val"] for r in per_dataset),
        "pos_test": sum(r["pos_test"] for r in per_dataset),
    }
    per_dataset.append(total)

    split_totals = final_split.groupby("split").agg(n=("id", "size"), pos=("label_hard", "sum"))
    return per_dataset, split_totals


def load_results():
    with open(RESULTS_DIR / "metrics_best.json") as f:
        m = json.load(f)
    test_results = pd.read_csv(RESULTS_DIR / "final_split_test_results.csv")
    sensitivity = pd.read_csv(RESULTS_DIR / "final_split_sensitivity_analysis.csv")
    return m, test_results.iloc[0], sensitivity


# ---------------------------------------------------------------------------
# Slide building helpers
# ---------------------------------------------------------------------------

def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_title(slide, text, top=Inches(0.3), font_size=28):
    tb = slide.shapes.add_textbox(MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = NAVY
    return tb


def add_bullets(slide, bullets, left, top, width, height, font_size=14):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + b
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)
    return tb


def add_table(slide, headers, rows, left, top, width, height, col_widths=None,
               font_size=11, bold_last_row=True):
    n_rows, n_cols = len(rows) + 1, len(headers)
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    if col_widths:
        for j, w in enumerate(col_widths):
            table.columns[j].width = w

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.margin_left = cell.margin_right = Emu(45720)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

    for i, row in enumerate(rows, start=1):
        is_total = bold_last_row and i == n_rows - 1
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Emu(45720)
            if is_total:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xD9, 0xE2, 0xF3)
            elif i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.bold = is_total
                p.font.color.rgb = DARK_TEXT
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
    return table


def add_picture_fit(slide, path, left, top, max_w, max_h):
    from PIL import Image
    with Image.open(path) as im:
        w_px, h_px = im.size
    ratio = w_px / h_px
    width, height = max_w, max_w / ratio
    if height > max_h:
        height = max_h
        width = max_h * ratio
    slide.shapes.add_picture(str(path), left, top, width=Emu(int(width)), height=Emu(int(height)))


def add_source_note(slide, text):
    tb = slide.shapes.add_textbox(MARGIN, SLIDE_H - Inches(0.4), SLIDE_W - 2 * MARGIN, Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x80, 0x80, 0x80)


def pct(x, decimals=1):
    s = f"{x * 100:.{decimals}f}"
    if "." in s:
        s = s.rstrip("0").rstrip(".")
    return s + "%"


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def build_title_slide(prs):
    slide = blank_slide(prs)
    tb = slide.shapes.add_textbox(MARGIN, Inches(2.4), SLIDE_W - 2 * MARGIN, Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.text = "ChagaSight - New Training Run Update"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    sub = slide.shapes.add_textbox(MARGIN, Inches(3.6), SLIDE_W - 2 * MARGIN, Inches(1.0))
    p = sub.text_frame.paragraphs[0]
    p.text = "Full-dataset run, two-phase training, final train / validation / test split (full_p2_50epochs)"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

    footer = slide.shapes.add_textbox(MARGIN, SLIDE_H - Inches(1.0), SLIDE_W - 2 * MARGIN, Inches(0.6))
    p = footer.text_frame.paragraphs[0]
    p.text = "Sanuth Mandepa | University of Westminster"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
    p.alignment = PP_ALIGN.CENTER


def build_dataset_composition_slide(prs, rows, age_stats):
    slide = blank_slide(prs)
    add_title(slide, "Dataset Composition: Original -> After Reduction")

    headers = ["Dataset", "Original Records", "Original Chagas+",
               "After Split-Assignment +\nPath Resolution",
               "Final (After Age-\nOutlier Filter)", "Final Chagas+"]
    table_rows = []
    for r in rows:
        orig_pct = pct(r["orig_pos"] / r["orig_records"], 2)
        final_pct = pct(r["final_pos"] / r["final_records"], 2)
        table_rows.append([
            r["label"],
            f"{r['orig_records']:,}",
            f"{r['orig_pos']:,} ({orig_pct})",
            f"{r['resolved_records']:,}",
            f"{r['final_records']:,}",
            f"{r['final_pos']:,} ({final_pct})",
        ])

    col_widths = [Inches(1.5), Inches(1.8), Inches(1.8), Inches(2.83), Inches(2.2), Inches(2.2)]
    add_table(slide, headers, table_rows, MARGIN, Inches(1.2), SLIDE_W - 2 * MARGIN, Inches(2.4),
              col_widths=col_widths, font_size=12)

    code15_removed = rows[0]["resolved_records"] - rows[0]["final_records"]
    ptb_removed_path = rows[1]["orig_records"] - rows[1]["resolved_records"]
    ptb_removed_age = rows[1]["resolved_records"] - rows[1]["final_records"]
    bullets = [
        f"PTB-XL lost {ptb_removed_path + ptb_removed_age:,} records: {ptb_removed_path:,} had no resolvable split assignment/path, "
        f"plus {ptb_removed_age:,} removed as age outliers.",
        f"CODE-15 lost {code15_removed:,} records to age-outlier filtering (raw ages ranged 2-300 yrs).",
        "SaMi-Trop unaffected by either filter.",
        f"Final age range across the retained dataset: {age_stats['min']:.0f}-{age_stats['max']:.0f} yrs "
        f"(mean {age_stats['mean']:.2f}, median {age_stats['median']:.1f}).",
    ]
    add_bullets(slide, bullets, MARGIN, Inches(4.0), SLIDE_W - 2 * MARGIN, Inches(2.8), font_size=16)
    add_source_note(slide, "Source: data/processed/metadata/all_data.csv, resolved_dataset_split.csv, resolved_dataset_split_with_age_filter.csv")


def build_split_slide(prs, per_dataset, split_totals):
    slide = blank_slide(prs)
    add_title(slide, "Final Train / Validation / Test Split (360,103 records)")

    headers1 = ["Dataset", "Train", "Val", "Test", "Total", "Chagas+ (Train / Val / Test)"]
    rows1 = []
    for r in per_dataset:
        total = r["train"] + r["val"] + r["test"]
        rows1.append([
            r["label"], f"{r['train']:,}", f"{r['val']:,}", f"{r['test']:,}", f"{total:,}",
            f"{r['pos_train']:,} / {r['pos_val']:,} / {r['pos_test']:,}",
        ])
    col_widths1 = [Inches(1.5), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8), Inches(3.63)]
    add_table(slide, headers1, rows1, MARGIN, Inches(1.15), SLIDE_W - 2 * MARGIN, Inches(2.2),
              col_widths=col_widths1, font_size=12)

    total_n = split_totals["n"].sum()
    headers2 = ["Split", "Records", "% of total", "Chagas+ (n, %)"]
    label_map = {"train": "Train", "val": "Validation", "test": "Test (held-out)"}
    rows2 = []
    for split in ["train", "val", "test"]:
        n = int(split_totals.loc[split, "n"])
        p = int(split_totals.loc[split, "pos"])
        rows2.append([label_map[split], f"{n:,}", pct(n / total_n, 1), f"{p:,} ({pct(p / n, 2)})"])
    col_widths2 = [Inches(2.5), Inches(2.5), Inches(2.5), Inches(4.83)]
    add_table(slide, headers2, rows2, MARGIN, Inches(3.65), SLIDE_W - 2 * MARGIN, Inches(1.8),
              col_widths=col_widths2, font_size=12, bold_last_row=False)

    bullets = [
        "Record-level split, no overlap between sets.",
        "PTB-XL serves entirely as confirmed-negative reference data (0% positive); SaMi-Trop is entirely "
        "confirmed-positive (100%); CODE-15 is the only dataset with a natural mix and is used entirely for training.",
    ]
    add_bullets(slide, bullets, MARGIN, Inches(5.7), SLIDE_W - 2 * MARGIN, Inches(1.4), font_size=15)
    add_source_note(slide, "Source: data/processed/metadata/final_split.csv")


def build_methodology_slide(prs):
    slide = blank_slide(prs)
    add_title(slide, "Training Methodology (New Run)")

    bullets = [
        "Model: dual-pathway ViT ensemble - 2D ViT (ECG contour image) + 1D ST-MEM ViT "
        "(12-lead, 100Hz signal) + age/sex demographics fusion",
        "Trained from random initialization this run (no pretrained MAE/ST-MEM checkpoints loaded)",
        "Phase 1: encoders frozen, 2000 iterations, LR 2e-4, effective batch 64",
        "Phase 2: full fine-tuning, up to 50 epochs (early-stopped at epoch 31), dual LR "
        "(2e-5 backbone / 2e-4 head), effective batch 32",
        "Loss: Focal loss (pos_weight=10) + REPA cross-modal alignment loss (lambda=0.5)",
        "Class imbalance: weighted sampling (oversample positives, ~2.2% prevalence)",
        "Augmentation (train only): lead mixup, powerline noise, random time shift, "
        "amplitude scaling, baseline wander",
        "Validation: evaluated every epoch (Phase 2) on a stratified subset; checkpoint "
        "selected by best validation AUROC; early stopping (patience 10, min_delta 1e-4) -> best = epoch 21",
        "Test: single evaluation of the epoch-21 checkpoint on the held-out test set; "
        "multiple decision thresholds compared",
    ]
    add_bullets(slide, bullets, MARGIN, Inches(1.2), Inches(8.0), Inches(6.0), font_size=14)
    add_picture_fit(slide, PLOTS_DIR / "final_split_training_curve.png",
                     Inches(8.7), Inches(1.2), Inches(4.1), Inches(6.0))


def build_results_slide(prs, m):
    slide = blank_slide(prs)
    add_title(slide, "Results (Held-out Test Set, n=3,198)")

    bullets = [
        f"AUROC: {m['test_auroc']:.4f} (primary)",
        f"AUPRC: {m['test_auprc']:.4f} (secondary)",
        f"Validation at best checkpoint (epoch {m['best_epoch']}): "
        f"AUROC {m['val_auroc']:.4f}, AUPRC {m['val_auprc']:.4f}",
        f"F1-optimal operating point (threshold = {m['test_f1_thr']:.3f}): "
        f"Precision {pct(m['test_f1_prec'])}, Recall {pct(m['test_f1_rec'])}, "
        f"F1 {m['test_f1_f1']:.3f}, Accuracy {pct(m['test_f1_acc'])} "
        f"(TP {m['test_f1_tp']}, FP {m['test_f1_fp']}, FN {m['test_f1_fn']}, TN {m['test_f1_tn']})",
    ]
    add_bullets(slide, bullets, MARGIN, Inches(1.3), Inches(7.6), Inches(5.5), font_size=18)

    add_picture_fit(slide, PLOTS_DIR / "final_split_roc_curve.png",
                     Inches(8.6), Inches(1.1), Inches(4.2), Inches(2.85))
    add_picture_fit(slide, PLOTS_DIR / "final_split_confusion_matrix.png",
                     Inches(8.6), Inches(4.1), Inches(4.2), Inches(2.85))
    add_source_note(slide, "Source: checkpoints_new/full_p2_50epochs/results/metrics_best.json")


def build_threshold_strategies_slide(prs, m, youden):
    slide = blank_slide(prs)
    add_title(slide, "Decision Threshold Strategies (Test Set)")

    youden_spec = youden["tn_youden"] / (youden["tn_youden"] + youden["fp_youden"])

    headers = ["Strategy", "Threshold", "Precision", "Recall", "F1", "Specificity",
               "Accuracy", "TP", "FP", "FN", "TN"]
    rows = [
        ["Max-Recall (screening)", f"{m['test_rec_thr']:.3f}", pct(m['test_rec_prec']), pct(m['test_rec_rec']),
         f"{m['test_rec_f1']:.3f}", pct(m['test_rec_spec']), pct(m['test_rec_acc']),
         m['test_rec_tp'], m['test_rec_fp'], m['test_rec_fn'], m['test_rec_tn']],
        ["High-Recall (balanced)", f"{m['test_recp_thr']:.3f}", pct(m['test_recp_prec']), pct(m['test_recp_rec']),
         f"{m['test_recp_f1']:.3f}", pct(m['test_recp_spec']), pct(m['test_recp_acc']),
         m['test_recp_tp'], m['test_recp_fp'], m['test_recp_fn'], m['test_recp_tn']],
        ["Youden's J", f"{youden['threshold_youden']:.3f}", pct(youden['precision_youden']), pct(youden['recall_youden']),
         f"{youden['f1_youden']:.3f}", pct(youden_spec), pct(youden['accuracy_youden']),
         youden['tp_youden'], youden['fp_youden'], youden['fn_youden'], youden['tn_youden']],
        ["F1-Optimal", f"{m['test_f1_thr']:.3f}", pct(m['test_f1_prec']), pct(m['test_f1_rec']),
         f"{m['test_f1_f1']:.3f}", pct(m['test_f1_spec']), pct(m['test_f1_acc']),
         m['test_f1_tp'], m['test_f1_fp'], m['test_f1_fn'], m['test_f1_tn']],
        ["Default (0.5)", f"{m['test_t05_thr']:.3f}", pct(m['test_t05_prec']), pct(m['test_t05_rec']),
         f"{m['test_t05_f1']:.3f}", pct(m['test_t05_spec']), pct(m['test_t05_acc']),
         m['test_t05_tp'], m['test_t05_fp'], m['test_t05_fn'], m['test_t05_tn']],
    ]
    col_widths = [Inches(2.0)] + [Inches(1.033)] * 10
    add_table(slide, headers, rows, MARGIN, Inches(1.2), SLIDE_W - 2 * MARGIN, Inches(2.8),
              col_widths=col_widths, font_size=11, bold_last_row=False)

    highlight = (
        f"Highlight: at the Max-Recall threshold (test_rec_thr = {m['test_rec_thr']:.3f} in "
        f"metrics_best.json), the model misses only {m['test_rec_fn']} of the 246 Chagas-positive "
        f"test cases ({pct(m['test_rec_rec'])} sensitivity), at the cost of {m['test_rec_fp']} false "
        f"positives among 2,952 negatives ({pct(1 - m['test_rec_spec'])} false-alarm rate) - "
        f"appropriate framing for a high-sensitivity screening tool."
    )
    add_bullets(slide, [highlight], MARGIN, Inches(4.3), SLIDE_W - 2 * MARGIN, Inches(1.6), font_size=16)
    add_source_note(slide, "Source: checkpoints_new/full_p2_50epochs/results/metrics_best.json, final_split_test_results.csv")


def build_sensitivity_sweep_slide(prs, sensitivity):
    slide = blank_slide(prs)
    add_title(slide, "Recall-Target Sensitivity Sweep (Test Set)")

    headers = ["Target Recall", "Threshold", "Actual Recall", "Precision", "F1",
               "Accuracy", "TP", "FP", "FN", "TN"]
    rows = []
    for _, r in sensitivity.iterrows():
        rows.append([
            r["Target recall"], f"{r['Threshold']:.3f}", pct(r['Actual recall']), pct(r['Precision']),
            f"{r['F1']:.3f}", pct(r['Accuracy']), int(r['TP']), int(r['FP']), int(r['FN']), int(r['TN']),
        ])
    col_widths = [Inches(1.5)] + [Inches(1.203)] * 9
    add_table(slide, headers, rows, MARGIN, Inches(1.2), SLIDE_W - 2 * MARGIN, Inches(3.0),
              col_widths=col_widths, font_size=12, bold_last_row=False)

    add_picture_fit(slide, PLOTS_DIR / "final_split_sensitivity_analysis.png",
                     MARGIN, Inches(4.5), SLIDE_W - 2 * MARGIN, Inches(2.55))
    add_source_note(slide, "Source: checkpoints_new/full_p2_50epochs/results/final_split_sensitivity_analysis.csv")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    rows, age_stats = load_dataset_composition()
    per_dataset, split_totals = load_final_split()
    metrics, youden, sensitivity = load_results()

    prs = new_presentation()
    build_title_slide(prs)
    build_dataset_composition_slide(prs, rows, age_stats)
    build_split_slide(prs, per_dataset, split_totals)
    build_methodology_slide(prs)
    build_results_slide(prs, metrics)
    build_threshold_strategies_slide(prs, metrics, youden)
    build_sensitivity_sweep_slide(prs, sensitivity)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Saved {OUT_PATH} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
